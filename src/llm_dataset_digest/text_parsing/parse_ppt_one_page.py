"""parse and save text of a PowerPoint."""
import os
import re

from pptx import Presentation

from llm_dataset_digest.text_parsing.iter_to_nonempty_table_cells import (
    iter_to_nonempty_table_cells,
)

# noqa
from llm_dataset_digest.text_parsing.save_to_json import save_to_json


def parse_ppt_one_page(  # noqa
    ppt_file: str,
    ppt_path: str,
    line_len_threshold: int = 4,
    text_len_threshold: int = 4,
) -> None:
    """This function parse and save text of a PowerPoint into json file.

    Args:
        ppt_file: file object.
        ppt_path: Path of the PowerPoint.
        line_len_threshold: threshold of line numbers. Only parse text greater than
        this number. Default = 4.
        text_len_threshold: threshold of text string lenght. Only parse text greater
        than this number. Default = 4.

    Returns:
        json files.
    """
    file_name = ppt_file.split(".")[0]
    folder_name = ppt_path + "\\" + file_name

    # create the folder to contain the json file outputs
    if not os.path.exists(folder_name):  # noqa
        os.mkdir(folder_name)  # noqa

    prs = Presentation(ppt_path + "\\" + ppt_file)
    output_file = ""

    for i, slide in enumerate(prs.slides):
        output_page = ""
        output_text_page = ""

        if slide.shapes:
            # sort the shapes by coord for each page
            shapes_coord = []
            for shape in slide.shapes:
                # find the coordinates of shape
                x = round(shape.left.inches, 2)
                y = round(shape.top.inches, 2)
                width = round(shape.width.inches, 2)
                height = round(shape.height.inches, 2)
                # append shape
                shapes_coord.append((x, y, width, height, shape))
            # sort the list of shapes
            shapes_coord_sort = sorted(
                shapes_coord, key=lambda element: (element[0], element[1])
            )

            # create empty dict for each page
            dict_shape = {
                "text": "",
                "text_type": "",
                "doc_name": file_name,
                "doc_year": re.match(r".*([1-3][0-9]{3})", file_name).group(1),
                "doc_quarter": file_name[-1],
                "page": str(i + 1),
            }

            for _j, item in enumerate(shapes_coord_sort):
                if i + 1 < 10:  # noqa
                    # add a zero before numbering
                    item_name = "pg" + str(0) + str(i + 1) + "_text" + "_" + file_name
                else:
                    item_name = "pg" + str(i + 1) + "_text" + "_" + file_name
                shape = item[4]

                # if the shape is a text frame with text
                if (
                    hasattr(shape, "text")
                    and len(str(shape.text).strip().split(" ")) > 1
                ):
                    # get the text from the text frame, and replace multiple newlines
                    # into one newline
                    shape_text_ = re.sub("\n+", "\n", shape.text.strip())
                    shape_text = shape_text_ + "\n" + "\n"

                    # get the font size (not robust)
                    text_frame = shape.text_frame
                    paragraph = text_frame.paragraphs[0]
                    for run in paragraph.runs:
                        font = run.font
                        try:
                            # get font size and bold check (not robust)
                            font_size = font.size.pt
                        except:  # noqa
                            font_size = 0

                    # get the length of string for each line in the text frame
                    shape_len = [
                        (len(line.strip().split(" ")))
                        for line in (shape_text.strip().split("\n"))
                    ]

                    # if length is greater than threshold
                    if (  # noqa
                        (max(shape_len) > text_len_threshold)
                        and (font_size > 8 or font_size == 0)  # noqa
                        and (shape_text not in output_file)
                    ):
                        output_page += shape_text
                        output_text_page += shape_text

                    # if length not greater than threshold, check if the text is a title
                    elif (  # noqa
                        font_size > 20 or font_size == 0  # noqa
                    ) and shape_text not in output_file:
                        output_page += shape_text
                        output_text_page += shape_text

                    elif shape_text[-2] in ".":
                        output_page += shape_text
                        output_text_page += shape_text

                # for parsing table
                elif shape.has_table:
                    tables_pg = []
                    tables = shape.table
                    tables_pg.append(tables)
                    for table in tables_pg:
                        table_text_ = iter_to_nonempty_table_cells(
                            table, text_len_threshold
                        )
                        table_text = "".join(table_text_)
                    shape_len = [
                        (len(line.strip().split(" ")))
                        for line in (table_text.strip().split("\n"))
                    ]

                    if (
                        len(shape_len) > line_len_threshold
                        or max(shape_len) > text_len_threshold
                    ) and table_text not in output_file:
                        output_page += table_text
                        output_text_page += table_text

        if output_text_page.strip():
            # store into json file
            dict_shape["text"] = output_text_page
            dict_shape["text_type"] = "text"
            save_to_json(path=folder_name, item=dict_shape, item_name=item_name)

            # store a file for quality check
            output_file += (
                "\n\n" + "---------- Page " + str(i + 1) + "------------" + "\n"
            )
            output_file += output_page
            os.chdir(ppt_path)
            with open(file_name + ".txt", "w", encoding="utf-8") as text_file:  # noqa
                text_file.write(str(output_file))
    return folder_name

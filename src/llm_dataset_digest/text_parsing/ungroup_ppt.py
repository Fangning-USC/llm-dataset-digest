"""Ungroup ppt function."""
import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def ungroup_ppt(
    input_file_path: str, output_file_name: str, output_file_path: str
) -> None:
    """This function ungroup all the grouped shapes in a PowerPoint and save it.

    Args:
        input_file_path: path of the directory.
        output_file_name: name of the ungrouped output ppt.
        output_file_path: path of the ungrouped output ppt.

    Returns:
        a PowerPoint file.
    """
    # ungroup all the grouped shapes so that the parser can recognize the text frames
    # within a group.
    presentation = Presentation(input_file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group = shape.element
                parent = group.getparent()
                index = parent.index(group)
                for member in group:
                    parent.insert(index, member)
                    index += 1
                parent.remove(group)
    try:
        # if output save path exist
        os.chdir(output_file_path)
        presentation.save(output_file_name)
    except:  # noqa
        # if output save path not exist, mkdir the output path
        os.mkdir(output_file_path)  # noqa
        os.chdir(output_file_path)  # noqa
        presentation.save(output_file_name)
